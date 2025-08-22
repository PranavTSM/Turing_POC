import toml
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE  # + MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.dml import MSO_THEME_COLOR
from pptx import Presentation
import csv
import json
import os

def _hex_to_rgb(hex_color: str):
    hex_color = hex_color.strip().lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join([c*2 for c in hex_color])
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)
def _soft_wrap_long_tokens(text: str, chunk=24):
    """
    Insert zero‑width spaces into very long tokens so PowerPoint can wrap them.
    e.g., 'superlongtoken...' -> 'superlongtoken<ZWSP>...'
    """
    parts = []
    for tok in str(text or "").split():
        if len(tok) > chunk:
            tok = "\u200b".join(tok[i:i+chunk] for i in range(0, len(tok), chunk))
        parts.append(tok)
    return " ".join(parts)


class Brand:
    def __init__(self, toml_file="config/brand.toml"):
        data = toml.load(toml_file) if os.path.exists(toml_file) else {}
        self.family_primary = data.get("typography", {}).get("family_primary", "Calibri")
        self.body_pt        = data.get("typography", {}).get("body_pt", 16)
        self.title_pt       = data.get("typography", {}).get("title_pt", 32)
        self.caption_pt     = data.get("typography", {}).get("caption_pt", 12)
        self.min_body_pt    = data.get("typography", {}).get("min_body_pt", 12)
        self.lh_title       = data.get("lineheights", {}).get("title", 1.15)
        self.lh_body        = data.get("lineheights", {}).get("body", 1.20)
        colors              = data.get("colors", {})
        self.color_text     = _hex_to_rgb(colors.get("text_primary", "#222222"))
        self.color_muted    = _hex_to_rgb(colors.get("text_secondary", "#4B4B4B"))
        self.color_accent   = _hex_to_rgb(colors.get("accent", "#005EB8"))
        self.table_header_bg= _hex_to_rgb(colors.get("table_header_bg", "#EEF5FD"))
        self.table_row_alt  = _hex_to_rgb(colors.get("table_row_alt", "#FAFAFA"))
        
    

class TemplateManager:
    def __init__(self, toml_file="config/templates.toml", brand_file="config/brand.toml"):
        self.templates = toml.load(toml_file)
        self.brand = Brand(brand_file)

    def get_template_names(self):
        return list(self.templates.keys())

    def get_placeholders(self, template_name):
        return self.templates.get(template_name, {}).get("placeholders", [])

    # ---------- helpers ----------
    def _apply_text(self, slide, ph, value):
        x, y, w, h = Inches(ph["x"]), Inches(ph["y"]), Inches(ph["w"]), Inches(ph["h"])
        tb = slide.shapes.add_textbox(x, y, w, h)

        # remove any visible border
        try:
            tb.line.fill.background()
            tb.line.width = 0
        except Exception:
            pass

        tf = tb.text_frame
        tf.clear()

        # keep box size fixed; let text wrap inside
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass

        # inner padding for nicer look
        tf.margin_left = tf.margin_right = Inches(0.12)
        tf.margin_top = tf.margin_bottom = Inches(0.12)

        # Vertical alignment
        valign = ph.get("valign", "top").lower()
        tf.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM
        }.get(valign, MSO_ANCHOR.TOP)

        p = tf.paragraphs[0]
        # Horizontal alignment
        align = ph.get("align", "left").lower()
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }.get(align, PP_ALIGN.LEFT)

        # Set text (with soft wrap for long tokens)
        text = _soft_wrap_long_tokens(value)
        p.text = str(text or "")

        # Font styling
        run = p.runs[0]
        font = run.font
        font.name = self.brand.family_primary
        base_size = ph.get("font_size", self.brand.body_pt)
        font.size = Pt(base_size)
        font.bold = bool(ph.get("bold", False))
        color = ph.get("color")
        font.color.rgb = _hex_to_rgb(color) if color else self.brand.color_text

        # Line spacing
        p.line_spacing = ph.get("line_spacing", self.brand.lh_body)

        # Optional shrink-to-min (heuristic)
        if ph.get("shrink_to_min", True):
            min_pt = self.brand.min_body_pt
            if len(p.text) > 600:
                size = base_size
                # shrink gradually for very long paragraphs
                while size > min_pt and len(p.text) > 600 + (base_size - size) * 120:
                    size -= 1
                    font.size = Pt(size)


    def _apply_image(self, slide, ph, path):
        if not path:
            return
        if not os.path.exists(path):
            print(f"⚠️ Image not found: {path}")
            return

        x, y, w, h = Inches(ph["x"]), Inches(ph["y"]), Inches(ph["w"]), Inches(ph["h"])
        pic = slide.shapes.add_picture(path, x, y, width=w, height=h)

        fit = (ph.get("fit") or "contain").lower()
        if fit == "contain":
            # already contained by width/height box
            return
        if fit == "cover":
            # Deterministic center-crop: compute aspect ratios, crop excess equally
            # python-pptx allows crop fractions 0..1
            # Here we approximate a centered crop by comparing image vs box ratio
            img_w, img_h = pic.image.size  # EMUs
            box_w, box_h = w, h  # EMUs (already in EMUs)
            ar_img = img_w / img_h
            ar_box = box_w / box_h

            if ar_img > ar_box:
                # too wide -> crop left/right
                new_w = ar_box * img_h
                excess = (img_w - new_w) / img_w
                left = right = excess / 2
                pic.crop_left = left
                pic.crop_right = left
            elif ar_img < ar_box:
                # too tall -> crop top/bottom
                new_h = img_w / ar_box
                excess = (img_h - new_h) / img_h
                top = bottom = excess / 2
                pic.crop_top = top
                pic.crop_bottom = top

    def _read_table_data(self, value):
        """
        Accepts:
          - path to CSV
          - JSON string '[[...], [...]]'
          - Python list of lists
        Returns (rows: List[List[str]])
        """
        if value is None:
            return []

        if isinstance(value, list):
            return value

        if isinstance(value, str) and os.path.exists(value) and value.lower().endswith(".csv"):
            rows = []
            with open(value, newline="", encoding="utf-8-sig") as f:
                for row in csv.reader(f):
                    rows.append(row)
            return rows

        # try JSON
        if isinstance(value, str):
            try:
                data = json.loads(value)
                if isinstance(data, list):
                    return data
            except Exception:
                pass

        # fallback single cell
        return [[str(value)]]

    def _apply_table(self, slide, ph, value):
        rows = self._read_table_data(value)
        if not rows:
            return

        x, y, w, h = Inches(ph["x"]), Inches(ph["y"]), Inches(ph["w"]), Inches(ph["h"])
        header = bool(ph.get("header", True))
        zebra  = bool(ph.get("zebra", True))
        font_size = Pt(ph.get("font_size", self.brand.body_pt))
        row_height_in = float(ph.get("row_height", 0.3))
        row_height = Inches(row_height_in)

        n_rows = len(rows)
        n_cols = len(rows[0])

        table_shape = slide.shapes.add_table(rows=n_rows, cols=n_cols, left=x, top=y, width=w, height=h)
        table = table_shape.table
        # table.alignment = WD_TABLE_ALIGNMENT.CENTER

        # Column widths (optional)
        col_widths = ph.get("columns")
        if col_widths:
            # Treat values as proportions or inches—either way normalize to the table width.
            widths = [float(v) for v in col_widths[:n_cols]]
            if len(widths) < n_cols:
                widths += [1.0] * (n_cols - len(widths))
            total = sum(widths) if sum(widths) > 0 else float(n_cols)
            for i in range(n_cols):
                frac = widths[i] / total
                table.columns[i].width = int(w * frac)  # EMU int required
        else:
            # equal widths
            for i in range(n_cols):
                table.columns[i].width = int(w / n_cols)


        # Fill cells
        # Fill cells
        for r in range(n_rows):
            for c in range(n_cols):
                cell = table.cell(r, c)

                # text with soft wrap for long tokens
                cell.text = _soft_wrap_long_tokens(rows[r][c])

                tf = cell.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = Inches(0.06)
                tf.margin_top = tf.margin_bottom = Inches(0.04)

                p = tf.paragraphs[0]
                p.font.size = font_size
                p.font.name = self.brand.family_primary
                p.font.color.rgb = self.brand.color_text
                p.alignment = PP_ALIGN.LEFT  # change per column 

                # reduce padding a bit
                cell.margin_left = cell.margin_right = cell.margin_top = cell.margin_bottom = 0

            # row height (approximate by adding newlines is unreliable; use row height property via xml)
            # python-pptx has no direct row height setter. We'll expand the overall table to fit:
        # Adjust table height to rows * row_height (+ header)
        desired_h = row_height * n_rows
        # lock width, set height
        table_shape.height = int(desired_h)


        # Header styling
        if header and n_rows >= 1:
            for c in range(n_cols):
                cell = table.cell(0, c)
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.color.rgb = self.brand.color_text
                # background
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.brand.table_header_bg

        # Zebra striping
        if zebra and n_rows > 2:
            for r in range(1 if header else 0, n_rows):
                if r % 2 == 1:
                    for c in range(n_cols):
                        cell = table.cell(r, c)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.brand.table_row_alt

    # ---------- main ----------
    def apply_content(self, slide, template_name, contents):
        placeholders = self.get_placeholders(template_name)

        for ph in placeholders:
            name = ph["name"]
            ph_type = ph["type"].lower()
            value = contents.get(name)

            if ph_type == "text":
                self._apply_text(slide, ph, value)
            elif ph_type == "image":
                self._apply_image(slide, ph, value)
            elif ph_type == "table":
                self._apply_table(slide, ph, value)
            else:
                print(f"⚠️ Unknown placeholder type: {ph_type}")

    # Convenience to create a slide by template name
    def add_slide(self, prs: Presentation, template_name: str, contents: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        self.apply_content(slide, template_name, contents)
        return slide
