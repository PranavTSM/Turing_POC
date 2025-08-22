# Keeps  public API but routes all work to TemplateManager so you get
# brand styling, text/images/tables, fit modes, etc.

from pptx import Presentation
from template import TemplateManager

def create_ppt_from_template(output_file, template, user_inputs):
    """
    Backward-compatible wrapper.

    Args:
        output_file (str): path for the .pptx to be written.
        template (str | dict): 
            - Preferred: a template *name* from config/templates.toml 
              (e.g., "Commentary_Only", "Charts_Grid", "Table_Only", ...).
            - Legacy: a dict with "placeholders" (old PoC). This path is
              deprecated; we’ll try to map it, but please migrate to a
              template name in templates.toml.
        user_inputs (dict): keys = placeholder names, values = text/image path/table data.
            - text: plain string
            - image: file path (png/jpg)
            - table: CSV path OR JSON of rows (e.g., [["H1","H2"],["R1C1","R1C2"]])
    """
    prs = Presentation()
    tm = TemplateManager("config/templates.toml", "config/brand.toml")

    # If caller passed a template name
    if isinstance(template, str) and template in tm.get_template_names():
        tm.add_slide(prs, template, user_inputs)

    # Legacy dict with "placeholders" (old PoC). We’ll synthesize a temporary slide.
    elif isinstance(template, dict) and "placeholders" in template:
        # Map legacy placeholder types/keys into our manager one-by-one.
        # NOTE: this path is limited; please migrate to TOML templates.
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for ph in template["placeholders"]:
            name = ph.get("name")
            ph_type = ph.get("type", "text").lower()
            # Legacy coords use x,y,cx,cy; new ones use x,y,w,h (in inches)
            x = ph.get("x") or ph.get("left") or ph.get("l") or 1
            y = ph.get("y") or ph.get("top")  or ph.get("t") or 1
            w = ph.get("w") or ph.get("cx")   or 4
            h = ph.get("h") or ph.get("cy")   or 1

            # Build a minimal ph dict compatible with TemplateManager internals
            normalized = {
                "name": name,
                "type": ph_type,
                "x": x, "y": y, "w": w, "h": h,
                "align": ph.get("alignment", ph.get("align", "left")),
                "valign": ph.get("valign", "top"),
                "font_size": ph.get("font_size"),
                "bold": ph.get("bold"),
                "fit": ph.get("fit", "contain"),
                "header": ph.get("header", True),
                "zebra": ph.get("zebra", True),
                "row_height": ph.get("row_height", 0.3),
                "columns": ph.get("columns"),
            }

            # Call the same private helpers the manager uses
            if ph_type == "text":
                tm._apply_text(slide, normalized, user_inputs.get(name, ""))
            elif ph_type == "image":
                tm._apply_image(slide, normalized, user_inputs.get(name))
            elif ph_type == "table":
                tm._apply_table(slide, normalized, user_inputs.get(name))
            else:
                print(f"⚠️ Unknown placeholder type in legacy template: {ph_type}")

    else:
        raise ValueError(
            "Unknown template reference. Pass a valid template name defined in "
            "config/templates.toml, or a legacy dict with 'placeholders'."
        )

    prs.save(output_file)
    print(f"✅ PPT created successfully: {output_file}")
