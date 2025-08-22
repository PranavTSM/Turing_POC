from pptx import Presentation
from template import TemplateManager

def _collect_inputs(placeholders):
    """
    Interactive prompt for values per placeholder.
    For 'table', you can paste a CSV path or JSON string like:
    [["Header1","Header2"],["R1C1","R1C2"]]
    """
    print("\n Provide content for placeholders:")
    data = {}
    for ph in placeholders:
        name = ph["name"]
        ph_type = ph["type"].lower()
        if ph_type == "text":
            value = input(f"Text for {name}: ")
        elif ph_type == "image":
            value = input(f"Image path for {name} (png/jpg): ")
        elif ph_type == "table":
            value = input(f"Table for {name} (CSV path or JSON rows): ")
        else:
            value = ""
        data[name] = value
    return data

def main():
    tm = TemplateManager("config/templates.toml", "config/brand.toml")
    prs = Presentation()

    print("Available templates:")
    for name in tm.get_template_names():
        print(" -", name)

    template_name = input("\nEnter template name: ").strip()
    if template_name not in tm.get_template_names():
        print(" Invalid template name.")
        return

    placeholders = tm.get_placeholders(template_name)
    sample_contents = _collect_inputs(placeholders)

    tm.add_slide(prs, template_name, sample_contents)

    out = "output.pptx"
    prs.save(out)
    print(f"✅ Presentation saved as {out}")

if __name__ == "__main__":
    main()
